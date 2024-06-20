from pptx import Presentation
from pptx.util import Inches
import os


def powerpoint_generate_for_julia(image_folder):
    prs = Presentation()
    image_dir = image_folder

    images = os.listdir(image_dir)


    image_groups = {}
    for image in images:
        if image.endswith(('.png', '.jpg')):
            #  I have two types of images naming as : Average of Log Transformed Intensity for ABL1 .png
            #                                         Log transformed Intensity for RAD9A for each DMSO and whel run.jpg
            if 'normalized Intensity difference' in image:
                gene_name = image.split('difference for ')[-1].split(' for each')[0]
                if gene_name not in image_groups:
                    image_groups[gene_name] = {}
                image_groups[gene_name]['difference'] = image
            elif 'normalized Intensity' in image:
                gene_name = image.split('Intensity for ')[-1].split(' for each')[0]
                if gene_name not in image_groups:
                    image_groups[gene_name] = {}
                image_groups[gene_name]['intensity'] = image
    # Got a dictionary of the set of files each with the average comparison and runs individual comparison.
    print(image_groups)


    for gene_name, images in image_groups.items():
        # Using a blank slide layout
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title_placeholder = slide.shapes.title
        title_placeholder.text = gene_name

        # Set the position of each image.
        if 'difference' in images:
            difference_image_path = os.path.join(image_dir, images['difference'])
            slide.shapes.add_picture(difference_image_path, Inches(0), Inches(2.5), width=Inches(5))
        if 'intensity' in images:
            intensity_image_path = os.path.join(image_dir, images['intensity'])
            slide.shapes.add_picture(intensity_image_path, Inches(5), Inches(2.75), width=Inches(5))

    prs.save('Powerpoint output/Z images powerpoint.pptx')

if __name__ == '__main__':
    powerpoint_generate_for_julia("powerpoint_img")